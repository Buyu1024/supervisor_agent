""".pptx PowerPoint 读取器 —— 递归形状遍历 + 表格/图表 + 空间排序

参考 MinerU 的 PPTX 处理管线：
1. 递归展开组合形状（GroupShape）
2. 提取表格、图表
3. XY-Cut 空间排序保持阅读顺序
4. 提取演讲者备注
"""

import tempfile
import os
from pathlib import Path
from pptx.enum.shapes import MSO_SHAPE_TYPE
from .base import BaseFileReader


class PptxReader(BaseFileReader):
    """
    读取 .pptx 文件，完整提取所有内容：

    处理能力：
    - 文本框（含嵌套组合形状内的）
    - 表格（行列单元格遍历）
    - 图表（标题 + 系列名）
    - 嵌入图片（OCR）
    - 演讲者备注
    - 空间排序（XY-Cut 近似阅读顺序）
    """

    suffixes = {'.pptx'}

    # 同一行的 Y 坐标容差（EMU 单位：1 inch = 914400 EMU）
    LINE_TOLERANCE = 50000  # ~0.05 inch

    def __init__(self):
        self._ocr = None  # 延迟加载

    def _get_ocr(self):
        """延迟初始化 RapidOCR"""
        if self._ocr is None:
            from rapidocr_onnxruntime import RapidOCR
            self._ocr = RapidOCR()
        return self._ocr

    # ================================================================
    # 主入口
    # ================================================================

    def read(self, path: Path) -> str:
        from pptx import Presentation

        prs = Presentation(str(path))
        slides_output: list[str] = []

        for slide_num, slide in enumerate(prs.slides, 1):
            # 递归收集当前幻灯片所有文本块（含位置信息）
            blocks = self._collect_blocks(slide.shapes)

            if not blocks:
                continue

            # XY-Cut 空间排序：先按 Y（行），再按 X（列）
            blocks.sort(key=lambda b: (b['y'], b['x']))

            slides_output.append(f"--- [幻灯片 {slide_num}] ---")
            for block in blocks:
                slides_output.append(block['text'])

            # 演讲者备注
            try:
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slides_output.append(f"[备注]: {notes}")
            except Exception:
                pass

        return '\n'.join(slides_output)

    # ================================================================
    # 递归形状遍历
    # ================================================================

    def _collect_blocks(
        self,
        shapes,
        parent_left: int = 0,
        parent_top: int = 0,
    ) -> list[dict]:
        """
        递归遍历形状树，收集所有文本块

        返回: [{"x": int, "y": int, "text": str}, ...]

        坐标转换（组合形状内的子形状需要加父偏移）：
        - parent_left / parent_top 是递归传递的累积偏移
        """
        blocks: list[dict] = []

        for shape in shapes:
            x = (shape.left or 0) + parent_left
            y = (shape.top or 0) + parent_top

            shape_type = shape.shape_type

            # ---- 组合形状：递归展开 ----
            if shape_type == MSO_SHAPE_TYPE.GROUP:
                blocks.extend(
                    self._collect_blocks(shape.shapes, x, y)
                )

            # ---- 表格 ----
            elif shape.has_table:
                blocks.extend(self._extract_table(shape, x, y))

            # ---- 图表 ----
            elif shape.has_chart:
                blocks.extend(self._extract_chart(shape, x, y))

            # ---- 文本框 ----
            elif shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    blocks.append({"x": x, "y": y, "text": text})

            # ---- 图片（OCR） ----
            if shape_type == MSO_SHAPE_TYPE.PICTURE:
                ocr_text = self._ocr_shape_image(shape)
                if ocr_text:
                    blocks.append({"x": x, "y": y, "text": ocr_text})

        return blocks

    # ================================================================
    # 表格提取
    # ================================================================

    def _extract_table(self, shape, base_x: int, base_y: int) -> list[dict]:
        """提取表格所有单元格文本"""
        blocks: list[dict] = []
        table = shape.table

        for row_idx, row in enumerate(table.rows):
            # 每行内部按列排序
            row_cells: list[dict] = []
            for col_idx, cell in enumerate(row.cells):
                text = cell.text.strip()
                if text:
                    row_cells.append({"col": col_idx, "text": text})

            if not row_cells:
                continue

            # 单列表格直接输出；多列用 | 拼接模拟表格
            if len(row_cells) == 1:
                blocks.append({
                    "x": base_x,
                    "y": base_y + row_idx * self.LINE_TOLERANCE,
                    "text": row_cells[0]["text"],
                })
            else:
                row_cells.sort(key=lambda c: c["col"])
                row_text = " | ".join(c["text"] for c in row_cells)
                blocks.append({
                    "x": base_x,
                    "y": base_y + row_idx * self.LINE_TOLERANCE,
                    "text": row_text,
                })

        return blocks

    # ================================================================
    # 图表提取
    # ================================================================

    def _extract_chart(self, shape, base_x: int, base_y: int) -> list[dict]:
        """提取图表标题和系列信息"""
        blocks: list[dict] = []
        chart = shape.chart
        offset = 0

        # 图表标题
        if chart.has_title:
            title_text = chart.chart_title.text_frame.text.strip()
            if title_text:
                blocks.append({
                    "x": base_x,
                    "y": base_y - self.LINE_TOLERANCE,
                    "text": title_text,
                })

        # 系列名称
        for series in chart.series:
            blocks.append({
                "x": base_x,
                "y": base_y + offset,
                "text": f"[图表系列] {series.name}",
            })
            offset += self.LINE_TOLERANCE

        return blocks

    # ================================================================
    # 嵌入图片 OCR
    # ================================================================

    def _ocr_shape_image(self, shape) -> str:
        """对形状内嵌的图片执行 OCR"""
        try:
            image = shape.image
            suffix = image.content_type.split('/')[-1] or 'png'
            with tempfile.NamedTemporaryFile(
                suffix=f'.{suffix}', delete=False
            ) as tmp:
                tmp.write(image.blob)
                tmp_path = tmp.name

            try:
                ocr = self._get_ocr()
                result, _ = ocr(tmp_path)
                if result:
                    lines = [r[1] for r in result if r[1].strip()]
                    if lines:
                        return '\n'.join(lines)
            finally:
                os.unlink(tmp_path)
        except Exception:
            pass
        return ""
